"""Generate Smart Factory Operations Center PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ────────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
ACCENT     = RGBColor(0x16, 0x21, 0x3E)   # mid navy
HIGHLIGHT  = RGBColor(0x0F, 0x3D, 0x66)   # steel blue
ORANGE     = RGBColor(0xE9, 0x4E, 0x1B)   # Itility orange-ish
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
YELLOW     = RGBColor(0xFF, 0xD7, 0x00)
GREEN      = RGBColor(0x2E, 0xCC, 0x71)
RED        = RGBColor(0xE7, 0x4C, 0x3C)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, colour=DARK_BG):
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colour


def add_rect(slide, left, top, width, height, fill_colour, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()   # no border
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_colour
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=24, bold=False, colour=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = colour
    run.font.italic = italic
    return txBox


def add_para(tf, text, font_size=18, bold=False, colour=WHITE,
             align=PP_ALIGN.LEFT, italic=False, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(font_size)
    run.font.bold   = bold
    run.font.color.rgb = colour
    run.font.italic = italic
    return p


def slide_header(slide, title, subtitle=None):
    """Top accent bar + title."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), ORANGE)
    add_textbox(slide, title,
                Inches(0.5), Inches(0.15), Inches(12.3), Inches(0.7),
                font_size=32, bold=True, colour=WHITE)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.5), Inches(0.82), Inches(12), Inches(0.4),
                    font_size=16, colour=LIGHT_GRAY, italic=True)


def bullet_box(slide, items, left, top, width, height,
               font_size=18, title=None, title_colour=YELLOW,
               bullet_colour=WHITE, title_size=20):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size  = Pt(title_size)
        run.font.bold  = True
        run.font.color.rgb = title_colour
    for item in items:
        p = tf.paragraphs[0] if (first and not title) else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = ("• " if not item.startswith("  ") else "") + item
        run.font.size  = Pt(font_size)
        run.font.bold  = False
        run.font.color.rgb = bullet_colour
    return txBox


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    slide = blank_slide(prs)
    fill_bg(slide, DARK_BG)

    # Orange accent bar (thick)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), ORANGE)
    # Bottom bar
    add_rect(slide, 0, SLIDE_H - Inches(0.12), SLIDE_W, Inches(0.12), ORANGE)

    # Left accent stripe
    add_rect(slide, 0, Inches(0.12), Inches(0.12), SLIDE_H - Inches(0.24), HIGHLIGHT)

    # Title
    add_textbox(slide, "Smart Factory Operations Center",
                Inches(0.5), Inches(1.8), Inches(12.3), Inches(1.3),
                font_size=44, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide, "Multi-Agent AI System for IPC Fleet Management",
                Inches(0.5), Inches(3.2), Inches(12.3), Inches(0.6),
                font_size=22, colour=LIGHT_GRAY, align=PP_ALIGN.CENTER, italic=True)

    # Details row
    add_textbox(slide, "LangGraph  ·  Claude (Anthropic)  ·  Streamlit  ·  LangFuse  ·  SQLite",
                Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.5),
                font_size=16, colour=ORANGE, align=PP_ALIGN.CENTER)

    add_textbox(slide, "Jelmer Blaas  |  Itility  |  May 2026",
                Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
                font_size=14, colour=LIGHT_GRAY, align=PP_ALIGN.CENTER)


def slide_problem(prs):
    slide = blank_slide(prs)
    fill_bg(slide, DARK_BG)
    slide_header(slide, "The Problem", "IPC fleet at scale — guesswork is expensive")

    # Context box
    add_rect(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.6), ACCENT)
    bullet_box(slide,
        ["Hundreds of IPCs across pizza lines in Europe",
         "Each IPC monitors cheese distribution with computer vision",
         "Initial rollout: over-provisioned hardware for safety",
         "Now: prove the system works → right-size the fleet"],
        Inches(0.65), Inches(1.55), Inches(5.5), Inches(2.5),
        font_size=17, title="Context", title_colour=ORANGE)

    # Questions box
    add_rect(slide, Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.6), HIGHLIGHT)
    bullet_box(slide,
        ["Which IPCs are consistently under-utilised?",
         "Which are approaching dangerous load levels?",
         "Which are healthy → no action needed?",
         "What actions should be prioritised fleet-wide?",
         "How do we avoid repeating rejected decisions?"],
        Inches(6.95), Inches(1.55), Inches(5.7), Inches(3.5),
        font_size=17, title="Key Questions", title_colour=YELLOW)

    # Constraint callout
    add_rect(slide, Inches(6.8), Inches(5.3), Inches(6.0), Inches(1.5), ORANGE)
    add_textbox(slide,
        "CONSTRAINT: All recommendations must be reviewed and approved by a human operator before any action is persisted.",
        Inches(6.95), Inches(5.35), Inches(5.7), Inches(1.3),
        font_size=15, bold=True, colour=WHITE)


def slide_architecture(prs):
    slide = blank_slide(prs)
    fill_bg(slide, DARK_BG)
    slide_header(slide, "Solution Architecture", "4-agent LangGraph supervisor pattern")

    # Central box — Orchestrator
    add_rect(slide, Inches(4.7), Inches(1.4), Inches(3.9), Inches(1.1), ORANGE)
    add_textbox(slide, "Orchestrator Agent",
                Inches(4.75), Inches(1.5), Inches(3.8), Inches(0.5),
                font_size=17, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "routes · HITL gate · session memory",
                Inches(4.75), Inches(1.85), Inches(3.8), Inches(0.35),
                font_size=12, colour=WHITE, align=PP_ALIGN.CENTER)

    # Three sub-agents
    boxes = [
        (Inches(0.4),  Inches(3.5), "Fleet Analyst",         "read-only sensor analysis",     GREEN),
        (Inches(4.7),  Inches(3.5), "Recommendation Engine", "urgency scores & proposals",    YELLOW),
        (Inches(9.0),  Inches(3.5), "Memory Manager",        "SQLite decisions & prefs",      HIGHLIGHT),
    ]
    for lft, tp, name, sub, col in boxes:
        add_rect(slide, lft, tp, Inches(3.5), Inches(1.0), col)
        add_textbox(slide, name, lft + Inches(0.1), tp + Inches(0.08),
                    Inches(3.3), Inches(0.45), font_size=16, bold=True,
                    colour=DARK_BG, align=PP_ALIGN.CENTER)
        add_textbox(slide, sub, lft + Inches(0.1), tp + Inches(0.5),
                    Inches(3.3), Inches(0.35), font_size=12,
                    colour=DARK_BG, align=PP_ALIGN.CENTER)

    # Data stores row
    stores = [
        (Inches(0.4),  Inches(5.2), "CSV Sensor Data",  "220k rows · pandas"),
        (Inches(4.7),  Inches(5.2), "classifier.pkl",   "sklearn DT (offline trained)"),
        (Inches(9.0),  Inches(5.2), "SQLite memory.db", "decisions · preferences"),
    ]
    for lft, tp, name, sub in stores:
        add_rect(slide, lft, tp, Inches(3.5), Inches(0.8), ACCENT)
        add_textbox(slide, name, lft + Inches(0.1), tp + Inches(0.05),
                    Inches(3.3), Inches(0.35), font_size=13, bold=True,
                    colour=LIGHT_GRAY, align=PP_ALIGN.CENTER)
        add_textbox(slide, sub, lft + Inches(0.1), tp + Inches(0.38),
                    Inches(3.3), Inches(0.3), font_size=11,
                    colour=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # Streamlit at top
    add_rect(slide, Inches(4.7), Inches(0.2), Inches(3.9), Inches(0.7), ACCENT)
    add_textbox(slide, "Streamlit Chat Interface",
                Inches(4.75), Inches(0.3), Inches(3.8), Inches(0.4),
                font_size=14, bold=True, colour=LIGHT_GRAY, align=PP_ALIGN.CENTER)


def slide_data(prs):
    slide = blank_slide(prs)
    fill_bg(slide, DARK_BG)
    slide_header(slide, "The Data", "EDA first — everything else follows from it")

    # Schema table header
    add_rect(slide, Inches(0.5), Inches(1.35), Inches(6.2), Inches(0.4), HIGHLIGHT)
    add_textbox(slide, "Column            Type     Description",
                Inches(0.6), Inches(1.38), Inches(6.0), Inches(0.35),
                font_size=13, bold=True, colour=WHITE)

    schema = [
        ("IPC",          "string",  "IPC identifier (e.g. ITLT4301)"),
        ("Data Factory", "int",     "Factory ID 1–5"),
        ("time",         "date",    "DD/MM/YYYY — 61 days (May–Jun 2021)"),
        ("AvgValue",     "float",   "Mean CPU usage (MHz)"),
        ("CpuMHz",       "mixed",   "Rated clock — some multi-CPU strings"),
        ("MetricId",     "string",  "Always CpuUsageMHz → drop at load"),
    ]
    for i, (col, typ, desc) in enumerate(schema):
        bg = ACCENT if i % 2 == 0 else DARK_BG
        add_rect(slide, Inches(0.5), Inches(1.75 + i*0.4), Inches(6.2), Inches(0.4), bg)
        add_textbox(slide, f"{col:<16} {typ:<8} {desc}",
                    Inches(0.6), Inches(1.78 + i*0.4), Inches(6.0), Inches(0.35),
                    font_size=12, colour=WHITE)

    # Key EDA findings
    add_rect(slide, Inches(7.0), Inches(1.35), Inches(5.8), Inches(5.6), ACCENT)
    bullet_box(slide,
        ["220,294 rows across 4,261 unique IPCs",
         "5 factories · 61-day date range",
         "AvgValue / CpuMHz × 100 = utilisation %",
         "Distribution is right-skewed → use p95, not mean",
         "AvgValue, MinValue, MaxValue are highly collinear (ρ≈0.9+) → use AvgValue only",
         "Factory 5 outlier: ITLT1593 reports >300,000 MHz on a 9,600 MHz CPU → excluded from training",
         "261 rows have composite CpuMHz strings → extract first numeric token"],
        Inches(7.1), Inches(1.5), Inches(5.5), Inches(5.2),
        font_size=15, title="Key EDA Findings", title_colour=YELLOW)


def slide_classifier(prs):
    slide = blank_slide(prs)
    fill_bg(slide, DARK_BG)
    slide_header(slide, "Classifier Pipeline", "Data-mined thresholds, not guesswork")

    steps = [
        ("1", "Aggregate",     "One row per IPC: cpu_avg, cpu_p95, cpu_max, days_observed",    HIGHLIGHT),
        ("2", "KMeans (k=4)",  "StandardScaler → cluster into 4 natural groups",               HIGHLIGHT),
        ("3", "Label Clusters","Inspect centroids → assign: underutilised / healthy / at-risk / overloaded", HIGHLIGHT),
        ("4", "Decision Tree", "DecisionTreeClassifier(max_depth=3) on cluster labels",         ORANGE),
        ("5", "Export",        "export_text() → classifier.pkl + scaler.pkl + tree_rules.txt",  ORANGE),
    ]
    for i, (num, title, desc, col) in enumerate(steps):
        y = Inches(1.35) + i * Inches(1.0)
        add_rect(slide, Inches(0.4), y, Inches(0.6), Inches(0.75), col)
        add_textbox(slide, num, Inches(0.4), y + Inches(0.15),
                    Inches(0.6), Inches(0.45), font_size=20, bold=True,
                    colour=WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, Inches(1.1), y, Inches(11.7), Inches(0.75), ACCENT)
        add_textbox(slide, title, Inches(1.2), y + Inches(0.05),
                    Inches(2.5), Inches(0.4), font_size=15, bold=True, colour=YELLOW)
        add_textbox(slide, desc, Inches(3.8), y + Inches(0.1),
                    Inches(8.8), Inches(0.55), font_size=14, colour=WHITE)

    # Labels box
    add_rect(slide, Inches(0.4), Inches(6.4), Inches(12.5), Inches(0.85), HIGHLIGHT)
    labels = [
        ("underutilised", "cpu_p95 < 30%", GREEN),
        ("healthy",       "30% ≤ cpu_p95 < 65%", YELLOW),
        ("at-risk",       "65% ≤ cpu_p95 < 85%", ORANGE),
        ("overloaded",    "cpu_p95 ≥ 85%", RED),
    ]
    x_start = Inches(0.6)
    for label, rule, col in labels:
        add_textbox(slide, f"● {label}", x_start, Inches(6.45),
                    Inches(1.8), Inches(0.3), font_size=13, bold=True, colour=col)
        add_textbox(slide, rule, x_start, Inches(6.75),
                    Inches(1.8), Inches(0.3), font_size=11, colour=LIGHT_GRAY)
        x_start += Inches(3.1)

    # Why note
    add_textbox(slide,
        "Why a decision tree over raw cluster labels? Fast, deterministic, and human-readable at inference time.",
        Inches(0.4), Inches(6.0), Inches(12.5), Inches(0.35),
        font_size=12, colour=LIGHT_GRAY, italic=True)


def slide_agents(prs):
    slide = blank_slide(prs)
    fill_bg(slide, DARK_BG)
    slide_header(slide, "Agent Design", "Separation of concerns — each agent owns exactly one responsibility")

    agents = [
        ("Orchestrator", ORANGE, [
            "Sole interface to the operator",
            "Routes to specialist agents",
            "Enforces HITL gate",
            "Injects long-term memory context",
            "Never touches data or DB directly",
        ]),
        ("Fleet Analyst", GREEN, [
            "Read-only sensor analysis",
            "compute_utilization_stats",
            "get_fleet_summary",
            "flag_anomalies",
            "classify_ipc (loads pickle)",
            "Returns facts, never recommendations",
        ]),
        ("Recommendation Engine", YELLOW, [
            "Receives Analyst output",
            "score_ipc_urgency (0–100)",
            "generate_rightsizing_plan",
            "rank_fleet_by_priority",
            "check_past_decisions before proposing",
        ]),
        ("Memory Manager", HIGHLIGHT, [
            "All SQLite reads & writes",
            "save_decision (only post-approval)",
            "load_past_decisions",
            "save_operator_preference",
            "get_session_context (session start)",
        ]),
    ]

    for i, (name, col, items) in enumerate(agents):
        x = Inches(0.35) + i * Inches(3.25)
        add_rect(slide, x, Inches(1.35), Inches(3.1), Inches(0.5), col)
        add_textbox(slide, name, x + Inches(0.1), Inches(1.4),
                    Inches(2.9), Inches(0.4), font_size=15, bold=True,
                    colour=DARK_BG, align=PP_ALIGN.CENTER)
        add_rect(slide, x, Inches(1.85), Inches(3.1), Inches(5.35), ACCENT)
        bullet_box(slide, items, x + Inches(0.1), Inches(1.95),
                   Inches(2.9), Inches(5.1), font_size=13)


def slide_hitl(prs):
    slide = blank_slide(prs)
    fill_bg(slide, DARK_BG)
    slide_header(slide, "Human-in-the-Loop (HITL)", "Enforced architecturally, not just by instruction")

    # Flow diagram
    flow = [
        (LIGHT_GRAY, "Operator",     '"Which IPCs are over-provisioned?"'),
        (HIGHLIGHT,  "Orchestrator", "Calls Fleet Analyst → classify_ipc for flagged IPCs\nCalls load_past_decisions for each IPC"),
        (ORANGE,     "Orchestrator", "Presents ranked proposals with:\n  • classification label\n  • CPU p95 value\n  • rule that fired\n  • recommended action"),
        (LIGHT_GRAY, "Operator",     '"Approve ITLT4301. Reject ITLT2034 — new line next month."'),
        (GREEN,      "Memory Mgr",   "save_decision(ITLT4301, status=\'approved\')\nsave_decision(ITLT2034, status=\'rejected\', note=\'new line next month\')"),
        (HIGHLIGHT,  "Next session", "load_past_decisions → ITLT2034 not re-suggested"),
    ]
    for i, (col, actor, text) in enumerate(flow):
        y = Inches(1.35) + i * Inches(0.95)
        add_rect(slide, Inches(0.3), y, Inches(1.8), Inches(0.8), col)
        add_textbox(slide, actor, Inches(0.35), y + Inches(0.2),
                    Inches(1.7), Inches(0.45), font_size=13, bold=True,
                    colour=DARK_BG, align=PP_ALIGN.CENTER)
        add_rect(slide, Inches(2.25), y, Inches(10.7), Inches(0.8), ACCENT)
        add_textbox(slide, text, Inches(2.35), y + Inches(0.08),
                    Inches(10.5), Inches(0.72), font_size=13, colour=WHITE)


def slide_memory(prs):
    slide = blank_slide(prs)
    fill_bg(slide, DARK_BG)
    slide_header(slide, "Memory System", "Short-term for the session · Long-term for operator decisions")

    # Short-term
    add_rect(slide, Inches(0.4), Inches(1.35), Inches(6.0), Inches(5.6), ACCENT)
    bullet_box(slide,
        ["LangGraph MessagesState",
         "Full conversation history in every turn",
         "Streamlit session_state holds message list",
         "Scoped per session — cleared on restart",
         "Simple, predictable, no summarisation overhead",
         "",
         "Trade-off: ConversationBufferMemory vs ConversationSummaryMemory",
         "  → Buffer chosen: bounded sessions, simpler, more auditable"],
        Inches(0.55), Inches(1.5), Inches(5.7), Inches(5.2),
        font_size=15, title="Short-Term Memory", title_colour=YELLOW)

    # Long-term
    add_rect(slide, Inches(6.9), Inches(1.35), Inches(6.0), Inches(5.6), ACCENT)
    bullet_box(slide,
        ["SQLite file at memory/operations.db",
         "Auto-created on first run",
         "Docker volume-mounted (persists restarts)",
         "",
         "decisions table:",
         "  id, ipc_id, action, status,",
         "  reason, operator_note, timestamp",
         "",
         "preferences table:",
         "  key, value, updated_at",
         "",
         "get_session_context() injects last 10",
         "decisions into Orchestrator system prompt"],
        Inches(7.05), Inches(1.5), Inches(5.7), Inches(5.2),
        font_size=14, title="Long-Term Memory (SQLite)", title_colour=YELLOW)


def slide_stack(prs):
    slide = blank_slide(prs)
    fill_bg(slide, DARK_BG)
    slide_header(slide, "Tech Stack & Key Decisions", "Every choice has a reason")

    rows = [
        ("Agent Framework",  "LangGraph",           "Explicit state machine; interrupt_before for HITL; built-in MessagesState"),
        ("LLM",              "Claude (Anthropic)",   "claude-sonnet-4-5 by default; configurable via ANTHROPIC_MODEL env var"),
        ("Classifier",       "sklearn KMeans + DT",  "Trained offline once; loaded as pickle at runtime — zero ML at query time"),
        ("Frontend",         "Streamlit",            "st.chat_input/message — production-grade chat UI in ~50 lines, no Node.js needed"),
        ("Persistence",      "SQLite",               "Zero config, file-based, ships with Python — sufficient for hundreds of IPCs"),
        ("Observability",    "LangFuse Cloud",       "CallbackHandler registered at graph level — all child spans auto-captured"),
        ("Packaging",        "Docker Compose",       "Single app service + SQLite volume mount; one command to run"),
        ("Data parsing",     "pandas",               "sep=';', decimal=',' — European CSV format; cpu_pct = AvgValue / CpuMHz × 100"),
    ]

    add_rect(slide, Inches(0.4), Inches(1.35), Inches(12.5), Inches(0.45), HIGHLIGHT)
    add_textbox(slide,
        "Layer                   Choice                   Rationale",
        Inches(0.5), Inches(1.38), Inches(12.3), Inches(0.38),
        font_size=13, bold=True, colour=WHITE)

    for i, (layer, choice, rationale) in enumerate(rows):
        bg = ACCENT if i % 2 == 0 else DARK_BG
        y = Inches(1.8) + i * Inches(0.62)
        add_rect(slide, Inches(0.4), y, Inches(12.5), Inches(0.6), bg)
        add_textbox(slide, layer,    Inches(0.5), y + Inches(0.12), Inches(2.3), Inches(0.4), font_size=13, bold=True, colour=YELLOW)
        add_textbox(slide, choice,   Inches(2.9), y + Inches(0.12), Inches(2.5), Inches(0.4), font_size=13, colour=ORANGE)
        add_textbox(slide, rationale, Inches(5.5), y + Inches(0.08), Inches(7.2), Inches(0.5), font_size=12, colour=WHITE)


def slide_observability(prs):
    slide = blank_slide(prs)
    fill_bg(slide, DARK_BG)
    slide_header(slide, "Observability with LangFuse", "Every LLM call, tool call, and agent span is traced")

    # Trace tree
    add_rect(slide, Inches(0.4), Inches(1.35), Inches(6.5), Inches(5.6), ACCENT)
    trace_lines = [
        ("Trace: session_{id}", 0, YELLOW, True),
        ("└── Span: orchestrator.run", 1, WHITE, False),
        ("      ├── Generation: orchestrator.llm_call", 2, LIGHT_GRAY, False),
        ("      ├── Event: tool.load_past_decisions", 2, LIGHT_GRAY, False),
        ("      ├── Span: fleet_analyst.run", 2, GREEN, False),
        ("      │     ├── Event: tool.get_fleet_summary", 3, LIGHT_GRAY, False),
        ("      │     └── Event: tool.classify_ipc × N", 3, LIGHT_GRAY, False),
        ("      ├── Span: recommendation_engine.run", 2, YELLOW, False),
        ("      │     ├── Event: tool.check_past_decisions × N", 3, LIGHT_GRAY, False),
        ("      │     └── Event: tool.generate_rightsizing_plan × N", 3, LIGHT_GRAY, False),
        ("      └── Span: memory_manager.run [if approved]", 2, ORANGE, False),
        ("            └── Event: tool.save_decision", 3, LIGHT_GRAY, False),
    ]
    txBox = slide.shapes.add_textbox(Inches(0.55), Inches(1.5), Inches(6.2), Inches(5.2))
    tf = txBox.text_frame
    tf.word_wrap = False
    first = True
    for line, indent, col, bold in trace_lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = line
        run.font.size = Pt(12)
        run.font.bold = bold
        run.font.color.rgb = col

    # Setup box
    add_rect(slide, Inches(7.2), Inches(1.35), Inches(5.8), Inches(5.6), HIGHLIGHT)
    bullet_box(slide,
        ["CallbackHandler registered at graph level",
         "All child spans auto-captured — no per-tool instrumentation",
         "",
         "env vars required:",
         "  LANGFUSE_PUBLIC_KEY=pk-lf-...",
         "  LANGFUSE_SECRET_KEY=sk-lf-...",
         "  LANGFUSE_HOST=https://cloud.langfuse.com",
         "",
         "Free tier is sufficient for this project",
         "",
         "Traces link directly to the session that produced them → easy debugging",
         "Failure is immediately localised to the failing agent span"],
        Inches(7.35), Inches(1.5), Inches(5.5), Inches(5.2),
        font_size=14, title="Setup & Benefits", title_colour=YELLOW)


def slide_tradeoffs(prs):
    slide = blank_slide(prs)
    fill_bg(slide, DARK_BG)
    slide_header(slide, "Design Decisions & Trade-offs", "")

    decisions = [
        ("4 agents vs 1",
         "Each agent's prompt is focused; Fleet Analyst callable independently; failure localised in LangFuse traces",
         "More wiring, more prompts to maintain"),
        ("Data-mined thresholds vs hardcoded",
         "Thresholds emerge from actual fleet behaviour; every recommendation shows the exact rule that fired → operator trust",
         "Requires offline training step; clusters need human labelling"),
        ("Streamlit vs React",
         "Production-grade chat UI in ~50 lines; no Node.js, no build step; ships with the Python container",
         "Less control over UI; not suitable for complex interactive dashboards"),
        ("Offline classifier (pickle) vs live ML",
         "Classification is microseconds, deterministic, and explainable at runtime",
         "Thresholds don't update automatically when new data arrives"),
        ("SQLite vs external DB",
         "File-based, zero config, ships in Python stdlib, persists via Docker volume",
         "Not suitable if multiple concurrent writers or very large decision history"),
    ]

    for i, (decision, pro, con) in enumerate(decisions):
        y = Inches(1.35) + i * Inches(1.15)
        add_rect(slide, Inches(0.4), y, Inches(3.5), Inches(1.0), ACCENT)
        add_textbox(slide, decision, Inches(0.5), y + Inches(0.15),
                    Inches(3.3), Inches(0.75), font_size=14, bold=True, colour=YELLOW)
        add_rect(slide, Inches(4.0), y, Inches(4.8), Inches(1.0), DARK_BG)
        add_textbox(slide, "✓ " + pro, Inches(4.1), y + Inches(0.1),
                    Inches(4.6), Inches(0.8), font_size=12, colour=GREEN)
        add_rect(slide, Inches(8.9), y, Inches(4.4), Inches(1.0), DARK_BG)
        add_textbox(slide, "△ " + con, Inches(9.0), y + Inches(0.1),
                    Inches(4.2), Inches(0.8), font_size=12, colour=ORANGE)

    # Column headers
    add_textbox(slide, "Decision", Inches(0.5), Inches(1.08), Inches(3.3), Inches(0.25),
                font_size=12, bold=True, colour=LIGHT_GRAY)
    add_textbox(slide, "Benefit", Inches(4.1), Inches(1.08), Inches(4.6), Inches(0.25),
                font_size=12, bold=True, colour=LIGHT_GRAY)
    add_textbox(slide, "Trade-off", Inches(9.0), Inches(1.08), Inches(4.2), Inches(0.25),
                font_size=12, bold=True, colour=LIGHT_GRAY)


def slide_demo(prs):
    slide = blank_slide(prs)
    fill_bg(slide, DARK_BG)
    slide_header(slide, "Live Demo", "")

    steps = [
        ("1", "Architecture",      "4 agents, who owns what and why"),
        ("2", "tree_rules.txt",    '"The data defined the thresholds — not us"'),
        ("3", "Broad query",       '"Which IPCs are over-provisioned?"'),
        ("4", "HITL in action",    "Approve one, defer one with a reason"),
        ("5", "New session",       "Ask the same question → agent references the deferral"),
        ("6", "LangFuse trace",    "Click through the session trace, show nested spans"),
        ("7", "Trade-offs",        "What we built and what we'd add with one more day"),
    ]

    for i, (num, title, desc) in enumerate(steps):
        y = Inches(1.4) + i * Inches(0.82)
        add_rect(slide, Inches(0.4), y, Inches(0.7), Inches(0.65), ORANGE)
        add_textbox(slide, num, Inches(0.4), y + Inches(0.1),
                    Inches(0.7), Inches(0.45), font_size=20, bold=True,
                    colour=WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, Inches(1.25), y, Inches(11.4), Inches(0.65), ACCENT)
        add_textbox(slide, title, Inches(1.35), y + Inches(0.1),
                    Inches(2.8), Inches(0.45), font_size=15, bold=True, colour=YELLOW)
        add_textbox(slide, desc, Inches(4.2), y + Inches(0.1),
                    Inches(8.3), Inches(0.5), font_size=14, colour=WHITE)

    add_textbox(slide,
        "Fallback: if LangGraph HITL is too complex → agent emits [NEEDS_APPROVAL] prefix, Streamlit holds until operator responds",
        Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.3),
        font_size=11, colour=LIGHT_GRAY, italic=True)


def slide_closing(prs):
    slide = blank_slide(prs)
    fill_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), ORANGE)
    add_rect(slide, 0, SLIDE_H - Inches(0.12), SLIDE_W, Inches(0.12), ORANGE)

    add_textbox(slide, "What We'd Add With One More Day",
                Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.7),
                font_size=28, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)

    next_steps = [
        ("Trend analysis",          "Detect IPCs degrading over time, not just snapshots"),
        ("Confidence scores",       "Distance to decision tree boundary on each classification"),
        ("Per-factory breakdowns",  "Fleet summary drilled down to factory and production line"),
        ("Multi-turn HITL",         '"Approve but defer 30 days" → operator can adjust timing'),
        ("Auto-retraining",         "Re-run classifier when fresh sensor data arrives"),
        ("pytest suite",            "Unit tests on compute_utilization_stats, classify_ipc, save_decision"),
        ("Streaming responses",     "Show agent thinking step-by-step in Streamlit UI"),
    ]
    for i, (title, desc) in enumerate(next_steps):
        col = i % 2
        row = i // 2
        x = Inches(0.4) + col * Inches(6.5)
        y = Inches(1.3) + row * Inches(1.0)
        add_rect(slide, x, y, Inches(6.2), Inches(0.85), ACCENT)
        add_textbox(slide, title, x + Inches(0.1), y + Inches(0.05),
                    Inches(6.0), Inches(0.35), font_size=14, bold=True, colour=ORANGE)
        add_textbox(slide, desc, x + Inches(0.1), y + Inches(0.4),
                    Inches(6.0), Inches(0.4), font_size=12, colour=WHITE)

    add_textbox(slide, "Thank you",
                Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.6),
                font_size=28, bold=True, colour=ORANGE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()

    slide_title(prs)           # ~0:30
    slide_problem(prs)         # ~1:00
    slide_architecture(prs)    # ~1:30
    slide_data(prs)            # ~1:00
    slide_classifier(prs)      # ~1:00
    slide_agents(prs)          # ~1:00
    slide_hitl(prs)            # ~1:00
    slide_memory(prs)          # ~0:45
    slide_stack(prs)           # ~0:45
    slide_observability(prs)   # ~0:30
    slide_tradeoffs(prs)       # ~0:45  (optional, use if time)
    slide_demo(prs)            # (demo notes slide)
    slide_closing(prs)         # ~0:30

    out = r"c:\Users\Jelmer\Documents\Itility\the-smart-factory\Smart_Factory_Operations_Center.pptx"
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
